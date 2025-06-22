import io
import re
import PyPDF2
from pptx import Presentation
import streamlit as st


def extract_text_from_pdf(pdf_file):
    """
    Extract text from a PDF file.

    Args:
        pdf_file: The uploaded PDF file object

    Returns:
        str: Extracted text from the PDF
    """
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""

        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text() + "\n\n"

        # Clean the text
        text = re.sub(r"\s+", " ", text)
        text = text.strip()

        return text
    except Exception as e:
        st.error(f"Lỗi khi trích xuất văn bản từ PDF: {str(e)}")
        raise


def extract_text_from_pptx(pptx_file):
    """
    Extract text from a PowerPoint file.

    Args:
        pptx_file: The uploaded PPTX file object

    Returns:
        str: Extracted text from the presentation
    """
    try:
        pptx_data = pptx_file.getvalue()
        presentation = Presentation(io.BytesIO(pptx_data))

        text = ""
        for slide in presentation.slides:
            slide_text = ""
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + " "
                except:
                    # Skip shapes that don't have text
                    pass
            text += slide_text.strip() + "\n\n"

        # Clean the text
        text = re.sub(r"\s+", " ", text)
        text = text.strip()

        return text
    except Exception as e:
        st.error(f"Lỗi khi trích xuất văn bản từ PowerPoint: {str(e)}")
        raise
